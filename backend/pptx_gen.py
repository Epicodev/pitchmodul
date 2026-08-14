"""
PowerPoint (.pptx) renderer.
Bygger samme pitch deck som HTML, men som editerbar .pptx-fil.

Bruger:
- 16:9 slide-størrelse (1920x1080 px → 13.333" x 7.5")
- Epico brand-farver direkte
- DM Sans som primær font (fallback til Arial hvis ikke installeret)
- Tekst er editerbar i PowerPoint
"""
import re
from io import BytesIO
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from slide_library import select_slides
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Brand-farver (matcher styles.css) ----------
RED = RGBColor(0x69, 0x0F, 0x23)
RASPBERRY = RGBColor(0xE0, 0x1E, 0x37)
BLACK_CURRANT = RGBColor(0x1B, 0x1B, 0x50)
KIWI = RGBColor(0x4C, 0xE1, 0x7F)
BLUEBERRY = RGBColor(0x4B, 0x64, 0xEA)
GREY = RGBColor(0x24, 0x21, 0x26)
LIGHT_GREY = RGBColor(0x91, 0x91, 0x99)
BEIGE = RGBColor(0xE4, 0xE1, 0xDC)
ALU_GREY = RGBColor(0xE5, 0xE5, 0xE5)
RAW_SILK = RGBColor(0xFF, 0xFC, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- Layout-konstanter ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

FONT_DISPLAY = "DM Sans"  # Fallback til Arial Bold når ikke installeret
FONT_BODY = "DM Sans"


# ============================================================
# HELPERS
# ============================================================


def _headline(ctx: Dict[str, Any], key: str, fallback_eyebrow: str, fallback_heading: str) -> tuple:
    """Hent AI-genereret overskrift til et kunde-slide.

    Overskriften er den samme som i HTML-decket — de to formater maa ikke sige
    forskellige ting. **Stjernerne** bevares og bliver til accent-farve inde i
    _add_text, saa fremhaevningen ogsaa overlever i PowerPoint.
    """
    h = (ctx.get("headlines") or {}).get(key) or {}
    eyebrow = (h.get("eyebrow") or "").strip() or fallback_eyebrow
    heading = (h.get("heading") or "").strip() or fallback_heading
    return eyebrow.upper(), heading


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Sæt baggrundsfarve på en slide."""
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, fill: RGBColor, line: Optional[RGBColor] = None):
    """Tilføj et rektangel med fyldfarve."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()  # Ingen border
    else:
        shape.line.color.rgb = line
    return shape



# ── Tekst-tilpasning ────────────────────────────────────────────────────────
# PowerPoint-tekstbokse har fast geometri. AI-genereret tekst varierer i længde,
# så en overskrift der passer for én kunde løber ud over kanten for den næste og
# lapper ind over indholdet nedenunder. HTML-versionen ombryder og flytter sig;
# her skal vi selv måle og skrue ned.
#
# Vi måler med Arial, som ligger metrisk tæt på DM Sans men er en anelse bredere.
# Det gør estimatet konservativt — vi skruer hellere en smule for meget ned end
# at lade to tekstblokke overlappe.

_MEASURE_FONTS = {
    False: "/System/Library/Fonts/Supplemental/Arial.ttf",
    True: "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
}
_REF_PT = 100  # måler én gang ved denne størrelse og skalerer lineært
_font_cache: Dict[bool, Any] = {}
_width_cache: Dict[tuple, float] = {}

# Hvor meget luft der er mellem linjerne, og hvor langt ned vi må skrue
_LINE_HEIGHT = 1.22
_MIN_SCALE = 0.62      # aldrig under 62% af den ønskede størrelse
_MIN_PT = 8


def _measure_font(bold: bool):
    if bold not in _font_cache:
        try:
            from PIL import ImageFont
            _font_cache[bold] = ImageFont.truetype(_MEASURE_FONTS[bold], _REF_PT)
        except Exception:
            _font_cache[bold] = None  # falder tilbage til tegnbredde-tilnærmelse
    return _font_cache[bold]


def _text_width_pt(text: str, pt: float, bold: bool) -> float:
    """Bredden af en tekststreng i punkter ved given skriftstørrelse."""
    if not text:
        return 0.0
    key = (text, bold)
    if key not in _width_cache:
        font = _measure_font(bold)
        if font is not None:
            _width_cache[key] = font.getbbox(text)[2] / _REF_PT
        else:
            # Ingen skrifttype at måle med — 0.52 em pr. tegn er et rimeligt
            # gennemsnit for en grotesk med blandet store og små bogstaver.
            _width_cache[key] = len(text) * 0.52
    return _width_cache[key] * pt


def _wrapped_line_count(text: str, pt: float, bold: bool, width_emu) -> int:
    """Hvor mange linjer teksten fylder når den ombrydes til den givne bredde."""
    if not text:
        return 0
    width_pt = Emu(int(width_emu)).pt
    lines = 0
    for para in text.split("\n"):
        words = para.split()
        if not words:
            lines += 1
            continue
        current = ""
        for word in words:
            trial = f"{current} {word}".strip()
            if _text_width_pt(trial, pt, bold) <= width_pt or not current:
                current = trial
            else:
                lines += 1
                current = word
        if current:
            lines += 1
    return lines


def _fit_font_size(text: str, requested_pt: int, bold: bool, width_emu, height_emu) -> int:
    """Find den største skriftstørrelse hvor teksten kan være i kassen.

    Skruer kun ned, aldrig op — layoutet er designet til den ønskede størrelse,
    og en overskrift der pludselig blev større ville se lige så forkert ud.
    """
    if not text or not height_emu:
        return requested_pt

    height_pt = Emu(int(height_emu)).pt
    floor = max(_MIN_PT, int(requested_pt * _MIN_SCALE))

    pt = requested_pt
    while pt > floor:
        lines = _wrapped_line_count(text, pt, bold, width_emu)
        if lines * pt * _LINE_HEIGHT <= height_pt:
            return pt
        pt -= 1
    return floor


def _clip(text: str, limit: int) -> str:
    """Afkort ved ordgraense med ellipse, aldrig midt i et ord.

    En haard afkortning gav slides med 'med en dedikeret Resource Manager s'.
    Naar teksten alligevel skal skaeres, skal det se ud som en beslutning.
    """
    text = str(text or "").strip()
    if len(text) <= limit:
        return text
    cut = text[:limit].rsplit(" ", 1)[0].rstrip(" ,.;:—–-")
    return f"{cut}…" if cut else text[:limit]


def _split_accent(text: str) -> List[tuple]:
    """Del tekst op i (stykke, er_fremhævet) ud fra **stjerne**-markering."""
    parts = []
    for i, chunk in enumerate(re.split(r"\*\*(.+?)\*\*", text or "")):
        if chunk:
            parts.append((chunk, i % 2 == 1))
    return parts or [(text or "", False)]


def _add_text(
    slide,
    x, y, w, h,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = GREY,
    font_name: str = FONT_BODY,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    letter_spacing: Optional[float] = None,
    accent: Optional[RGBColor] = None,
    shrink: bool = True,
):
    """Tilføj en tekstboks der tilpasser sig indholdet.

    `shrink` skruer skriftstørrelsen ned indtil teksten kan være i kassen. Uden
    det ville lange AI-genererede overskrifter lappe ind over indholdet nedenunder.
    Sæt den til False for tekst hvor størrelsen er vigtigere end at alt er synligt.

    `accent` farver ord markeret med **stjerner** — samme fremhævning som i
    HTML-versionen, så de to formater viser den samme overskrift.
    """
    text = text or ""
    segments = _split_accent(text) if accent else [(text, False)]
    plain = "".join(seg for seg, _ in segments)

    if shrink:
        font_size = _fit_font_size(plain, font_size, bold, w, h)

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    for segment, is_accent in segments:
        if not segment:
            continue
        run = p.add_run()
        run.text = segment
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = accent if (is_accent and accent) else color
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))  # 1/100 pt

    return tb


def _add_e_mark(slide, x, y, size: int = Inches(0.5), bg: RGBColor = RED, fg: RGBColor = WHITE):
    """Tilføj Epico E-bomærket som en gruppe af rektangler."""
    # Baggrunds-firkant
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()

    # 3 horisontale streger der danner E
    bar_height = size // 7
    side_padding = size // 6
    bar_width_full = size - (side_padding * 2)
    bar_width_mid = int(bar_width_full * 0.65)
    bar_x = x + side_padding

    # Top bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, y + side_padding, bar_width_full, bar_height)
    top.fill.solid(); top.fill.fore_color.rgb = fg; top.line.fill.background()

    # Mid bar (kortere)
    mid_y = y + (size // 2) - (bar_height // 2)
    mid = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, mid_y, bar_width_mid, bar_height)
    mid.fill.solid(); mid.fill.fore_color.rgb = fg; mid.line.fill.background()

    # Bottom bar
    bot_y = y + size - side_padding - bar_height
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bot_y, bar_width_full, bar_height)
    bot.fill.solid(); bot.fill.fore_color.rgb = fg; bot.line.fill.background()


def _add_slide_header(slide, section_tag: str, dark: bool = False):
    """Tilføj brand-header i top-right (E-mærke + 'Epico' + section-tag)."""
    text_color = WHITE if dark else BLACK_CURRANT
    label_color = RGBColor(0xFF, 0xFF, 0xFF) if dark else LIGHT_GREY
    e_size = Inches(0.4)

    # E-mærke
    e_x = SLIDE_W - Inches(3.5)
    e_y = Inches(0.35)
    _add_e_mark(slide, e_x, e_y, e_size)

    # "Epico" wordmark
    _add_text(
        slide,
        e_x + e_size + Inches(0.15), e_y - Inches(0.02),
        Inches(0.8), e_size + Inches(0.1),
        "Epico",
        font_size=14, bold=True, color=text_color, anchor=MSO_ANCHOR.MIDDLE,
        font_name=FONT_DISPLAY,
    )

    # Section-tag
    _add_text(
        slide,
        e_x + e_size + Inches(1.0), e_y,
        Inches(2.0), e_size,
        section_tag.upper(),
        font_size=9, bold=True, color=label_color, anchor=MSO_ANCHOR.MIDDLE,
        font_name=FONT_BODY,
    )


def _add_footer(slide, text: str, dark: bool = False):
    """Tilføj slide-footer i bottom-left."""
    color = RGBColor(0xFF, 0xFF, 0xFF) if dark else LIGHT_GREY
    _add_text(
        slide,
        MARGIN, SLIDE_H - Inches(0.5),
        Inches(6), Inches(0.3),
        text.upper(),
        font_size=9, bold=True, color=color,
        font_name=FONT_BODY,
    )


# ============================================================
# SLIDE BUILDERS
# ============================================================

def _slide_cover(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, BLACK_CURRANT)

    # E-mærke + wordmark i top-left
    _add_e_mark(slide, MARGIN, Inches(0.6), Inches(0.55))
    _add_text(slide, MARGIN + Inches(0.7), Inches(0.6),
              Inches(2), Inches(0.55), "Epico",
              font_size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
              font_name=FONT_DISPLAY)

    # Meta top-right
    _add_text(slide, SLIDE_W - Inches(4.5), Inches(0.65),
              Inches(4), Inches(0.4),
              f"SKRÆDDERSYET PITCH · {ctx['meeting']['date'].upper()}",
              font_size=9, bold=True,
              color=RGBColor(0x99, 0x99, 0xAA), align=PP_ALIGN.RIGHT,
              font_name=FONT_BODY)

    # Eyebrow
    _add_text(slide, MARGIN, Inches(2.8),
              Inches(6), Inches(0.4),
              "STRATEGISK SAMARBEJDSOPLÆG",
              font_size=11, bold=True, color=KIWI,
              font_name=FONT_BODY)

    # Stort title — kunde-navn
    _add_text(slide, MARGIN, Inches(3.3),
              SLIDE_W - (MARGIN * 2), Inches(1.0),
              ctx["client"]["name"],
              font_size=72, bold=True, color=WHITE,
              font_name=FONT_DISPLAY)

    # X-mark
    _add_text(slide, MARGIN, Inches(4.4),
              Inches(2), Inches(0.6),
              "×",
              font_size=44, bold=True, color=KIWI,
              font_name=FONT_DISPLAY)

    # Epico
    _add_text(slide, MARGIN, Inches(5.0),
              SLIDE_W - (MARGIN * 2), Inches(1.0),
              "Epico",
              font_size=72, bold=True, color=KIWI,
              font_name=FONT_DISPLAY)

    # Footer
    _add_text(slide, MARGIN, SLIDE_H - Inches(0.7),
              Inches(6), Inches(0.3),
              f"UDARBEJDET TIL {ctx['meeting']['contact_person'].upper()}",
              font_size=9, bold=True,
              color=RGBColor(0xAA, 0xAA, 0xBB),
              font_name=FONT_BODY)
    _add_text(slide, SLIDE_W - Inches(6.5), SLIDE_H - Inches(0.7),
              Inches(6), Inches(0.3),
              f"{ctx['meeting']['city'].upper()} · {ctx['meeting']['date'].upper()}",
              font_size=9, bold=True,
              color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.RIGHT,
              font_name=FONT_BODY)


def _slide_agenda(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RAW_SILK)
    _add_slide_header(slide, "Agenda")

    _add_text(slide, MARGIN, Inches(1.4), Inches(6), Inches(0.4),
              "SÅDAN BRUGER VI DE NÆSTE 45 MINUTTER",
              font_size=11, bold=True, color=RASPBERRY,
              font_name=FONT_BODY)

    _add_text(slide, MARGIN, Inches(1.85), Inches(8), Inches(1.0),
              "Agenda",
              font_size=60, bold=True, color=BLACK_CURRANT,
              font_name=FONT_DISPLAY)

    items = [
        ("01", f"Vores research om {ctx['client']['name']}", "8 min"),
        ("02", "Jeres strategiske prioriteter — sådan som vi læser dem", "10 min"),
        ("03", "Hvor Epico kan flytte nålen for jer", "10 min"),
        ("04", "Hvem vi er, og hvad vi leverer", "10 min"),
        ("05", "Næste skridt — hvis vi er enige", "7 min"),
    ]
    y = Inches(3.2)
    for num, label, dur in items:
        _add_text(slide, MARGIN, y, Inches(0.7), Inches(0.5),
                  num, font_size=22, bold=True, color=RASPBERRY,
                  font_name=FONT_DISPLAY)
        _add_text(slide, MARGIN + Inches(0.9), y, Inches(8.5), Inches(0.5),
                  label, font_size=22, color=BLACK_CURRANT,
                  font_name=FONT_DISPLAY)
        _add_text(slide, SLIDE_W - Inches(1.5), y, Inches(1.0), Inches(0.5),
                  dur.upper(), font_size=10, bold=True, color=LIGHT_GREY,
                  align=PP_ALIGN.RIGHT, font_name=FONT_BODY)
        # Divider line
        _add_rect(slide, MARGIN, y + Inches(0.65), SLIDE_W - (MARGIN * 2), Emu(8000), ALU_GREY)
        y += Inches(0.78)


def _slide_research(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RAW_SILK)
    _add_slide_header(slide, "Research")

    eyebrow, heading = _headline(ctx, "research", "Vi har gjort hjemmearbejdet",
                                 f"Dette ved vi om {ctx['client']['name']}")
    _add_text(slide, MARGIN, Inches(1.4), Inches(6), Inches(0.4),
              eyebrow, font_size=11, bold=True, color=RASPBERRY, font_name=FONT_BODY)
    _add_text(slide, MARGIN, Inches(1.85), Inches(11), Inches(1.0),
              heading, font_size=56, bold=True, color=BLACK_CURRANT,
              font_name=FONT_DISPLAY, accent=RASPBERRY)

    # 2x2 grid med facts
    facts = ctx.get("research_facts", [])[:4]
    cell_w = (SLIDE_W - (MARGIN * 2) - Inches(0.05)) / 2
    cell_h = Inches(1.6)
    start_y = Inches(3.4)
    for i, fact in enumerate(facts):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (cell_w + Inches(0.05))
        y = start_y + row * (cell_h + Inches(0.05))
        _add_rect(slide, x, y, cell_w, cell_h, WHITE)
        # Venstre rød border
        _add_rect(slide, x, y, Inches(0.05), cell_h, RED)
        # Key
        _add_text(slide, x + Inches(0.4), y + Inches(0.25),
                  cell_w - Inches(0.8), Inches(0.3),
                  fact.get("key", "").upper(),
                  font_size=10, bold=True, color=LIGHT_GREY,
                  font_name=FONT_BODY)
        # Value
        value_text = fact.get("value", "")
        value_size = 28 if len(value_text) <= 18 else 18
        _add_text(slide, x + Inches(0.4), y + Inches(0.6),
                  cell_w - Inches(0.8), Inches(0.55),
                  value_text,
                  font_size=value_size, bold=True, color=BLACK_CURRANT,
                  font_name=FONT_DISPLAY)
        # Source
        _add_text(slide, x + Inches(0.4), y + cell_h - Inches(0.4),
                  cell_w - Inches(0.8), Inches(0.3),
                  f"Kilde: {fact.get('source', '')}",
                  font_size=10, color=LIGHT_GREY, font_name=FONT_BODY)


def _slide_mapping(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RAW_SILK)
    _add_slide_header(slide, "Hvor vi flytter nålen")

    eyebrow, heading = _headline(ctx, "mapping", "Konkret kobling",
                                 "Jeres udfordring. Vores håndtag.")
    _add_text(slide, MARGIN, Inches(1.4), Inches(6), Inches(0.4),
              eyebrow, font_size=11, bold=True, color=RASPBERRY, font_name=FONT_BODY)
    _add_text(slide, MARGIN, Inches(1.85), Inches(12), Inches(1.0),
              heading, font_size=44, bold=True, color=BLACK_CURRANT,
              font_name=FONT_DISPLAY, accent=RASPBERRY)

    mappings = ctx.get("value_mappings", [])[:4]
    col1_w = Inches(5.5)
    arrow_w = Inches(0.6)
    col2_w = SLIDE_W - (MARGIN * 2) - col1_w - arrow_w

    # Header row
    header_y = Inches(3.3)
    _add_rect(slide, MARGIN, header_y, col1_w, Inches(0.45), BEIGE)
    _add_text(slide, MARGIN + Inches(0.2), header_y + Inches(0.1),
              col1_w - Inches(0.3), Inches(0.3),
              "JERES UDFORDRING",
              font_size=10, bold=True, color=GREY, font_name=FONT_BODY)
    _add_rect(slide, MARGIN + col1_w + arrow_w, header_y, col2_w, Inches(0.45), BLACK_CURRANT)
    _add_text(slide, MARGIN + col1_w + arrow_w + Inches(0.2), header_y + Inches(0.1),
              col2_w - Inches(0.3), Inches(0.3),
              "DET EPICO KAN LØSE",
              font_size=10, bold=True, color=WHITE, font_name=FONT_BODY)

    y = header_y + Inches(0.5)
    row_h = Inches(0.75)
    for m in mappings:
        _add_rect(slide, MARGIN, y, col1_w, row_h, WHITE, line=ALU_GREY)
        _add_text(slide, MARGIN + Inches(0.2), y + Inches(0.1),
                  col1_w - Inches(0.4), row_h - Inches(0.2),
                  _clip(m.get("challenge", ""), 200),
                  font_size=11, color=GREY, font_name=FONT_BODY,
                  anchor=MSO_ANCHOR.MIDDLE)
        # Arrow
        _add_text(slide, MARGIN + col1_w, y,
                  arrow_w, row_h,
                  "→", font_size=20, bold=True, color=RASPBERRY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  font_name=FONT_DISPLAY)
        _add_rect(slide, MARGIN + col1_w + arrow_w, y, col2_w, row_h,
                  RGBColor(0xFA, 0xF9, 0xF5), line=ALU_GREY)
        sol_text = f"{m.get('epico_service', '')}: {_clip(m.get('solution', ''), 180)}"
        _add_text(slide, MARGIN + col1_w + arrow_w + Inches(0.2),
                  y + Inches(0.1),
                  col2_w - Inches(0.4), row_h - Inches(0.2),
                  sol_text,
                  font_size=11, color=BLACK_CURRANT, font_name=FONT_BODY,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + Emu(20000)


def _slide_divider(prs, chapter_num: str, title: str, accent_word: str = "", red: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RED if red else BLACK_CURRANT)
    _add_slide_header(slide, f"Kapitel {chapter_num}", dark=True)

    # Stort baggrunds-nummer
    _add_text(slide, MARGIN, Inches(1.0), Inches(6), Inches(4.5),
              chapter_num,
              font_size=240, bold=True,
              color=RGBColor(0x33, 0x33, 0x77) if not red else RGBColor(0x80, 0x20, 0x35),
              font_name=FONT_DISPLAY)

    _add_text(slide, MARGIN, Inches(5.8), Inches(6), Inches(0.4),
              f"KAPITEL {chapter_num}".upper(),
              font_size=11, bold=True, color=KIWI, font_name=FONT_BODY)

    _add_text(slide, MARGIN, Inches(6.2), SLIDE_W - (MARGIN * 2), Inches(1.0),
              title,
              font_size=80, bold=True, color=WHITE, font_name=FONT_DISPLAY)






def _slide_case(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RAW_SILK)
    _add_slide_header(slide, "Relevant case")

    case = ctx.get("case", {})
    _add_text(slide, MARGIN, Inches(1.0), Inches(8), Inches(0.4),
              f"CASE FRA {ctx.get('industry_tag', 'BRANCHEN').upper()}",
              font_size=11, bold=True, color=RASPBERRY,
              font_name=FONT_BODY)
    _add_text(slide, MARGIN, Inches(1.45), Inches(12), Inches(1.2),
              _clip(case.get("headline", ""), 180),
              font_size=36, bold=True, color=BLACK_CURRANT,
              font_name=FONT_DISPLAY)
    _add_text(slide, MARGIN, Inches(2.7), Inches(12), Inches(0.8),
              _clip(case.get("intro", ""), 300),
              font_size=14, color=GREY, font_name=FONT_BODY)

    # 4 kolonner
    cols = [
        ("HVAD", case.get("what", []), WHITE, GREY, RASPBERRY),
        ("HVORFOR", case.get("why", []), WHITE, GREY, RASPBERRY),
        ("RESULTAT", case.get("result", []), RED, WHITE, KIWI),
        ("VÆRDI", case.get("value", []), WHITE, GREY, RASPBERRY),
    ]
    # Case-bullets er hele saetninger fra AI'en. Med 0,55" pr. bullet skrumpede
    # de ned i 8pt for at passe — mens der stod en tom tomme i bunden af sliden.
    # Her faar de pladsen i stedet, saa teksten kan blive laesbar paa en projektor.
    col_w = (SLIDE_W - (MARGIN * 2) - Inches(0.15)) / 4
    cell_h = Inches(3.5)
    y = Inches(3.6)
    bullet_h = Inches(0.78)
    bullet_step = Inches(0.82)
    for i, (label, bullets, bg, txt, accent) in enumerate(cols):
        x = MARGIN + i * (col_w + Inches(0.05))
        _add_rect(slide, x, y, col_w, cell_h, bg)
        _add_text(slide, x + Inches(0.25), y + Inches(0.2),
                  col_w - Inches(0.5), Inches(0.3),
                  label, font_size=10, bold=True, color=accent,
                  font_name=FONT_BODY)
        bl_y = y + Inches(0.65)
        for b in bullets[:3]:
            _add_text(slide, x + Inches(0.25), bl_y,
                      col_w - Inches(0.5), bullet_h,
                      f"• {_clip(b, 150)}", font_size=11, color=txt,
                      font_name=FONT_BODY)
            bl_y += bullet_step


def _slide_next_steps(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, RAW_SILK)
    _add_slide_header(slide, "Næste skridt")

    eyebrow, heading = _headline(ctx, "next_steps", "Hvis vi er enige om retningen",
                                 "Tre konkrete næste skridt.")
    _add_text(slide, MARGIN, Inches(1.4), Inches(8), Inches(0.4),
              eyebrow, font_size=11, bold=True, color=RASPBERRY,
              font_name=FONT_BODY)
    _add_text(slide, MARGIN, Inches(1.85), Inches(11), Inches(1.0),
              heading, font_size=44, bold=True, color=BLACK_CURRANT,
              font_name=FONT_DISPLAY, accent=RASPBERRY)

    steps = ctx.get("next_steps", [])[:3]
    col_w = (SLIDE_W - (MARGIN * 2) - Inches(0.1)) / 3
    cell_h = Inches(3.0)
    y = Inches(3.4)
    for i, step in enumerate(steps, start=1):
        x = MARGIN + (i - 1) * (col_w + Inches(0.05))
        _add_rect(slide, x, y, col_w, cell_h, WHITE)
        _add_text(slide, x + Inches(0.3), y + Inches(0.3),
                  col_w - Inches(0.5), Inches(0.3),
                  f"SKRIDT {i:02d}", font_size=10, bold=True, color=RASPBERRY,
                  font_name=FONT_BODY)
        _add_text(slide, x + Inches(0.3), y + Inches(0.75),
                  col_w - Inches(0.5), Inches(0.8),
                  step.get("title", ""), font_size=20, bold=True,
                  color=BLACK_CURRANT, font_name=FONT_DISPLAY)
        _add_text(slide, x + Inches(0.3), y + Inches(1.65),
                  col_w - Inches(0.5), Inches(0.9),
                  _clip(step.get("description", ""), 200),
                  font_size=11, color=GREY, font_name=FONT_BODY)
        # When-badge
        _add_rect(slide, x + Inches(0.3), y + cell_h - Inches(0.55),
                  col_w - Inches(0.6), Emu(15000), ALU_GREY)
        _add_text(slide, x + Inches(0.3), y + cell_h - Inches(0.45),
                  col_w - Inches(0.6), Inches(0.3),
                  step.get("when", "").upper(), font_size=10, bold=True,
                  color=BLACK_CURRANT, font_name=FONT_BODY)


def _slide_contact(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BLACK_CURRANT)
    _add_slide_header(slide, "Kontakt", dark=True)

    _add_text(slide, MARGIN, Inches(1.5), Inches(8), Inches(0.4),
              "LAD OS TAGE NÆSTE SKRIDT",
              font_size=11, bold=True, color=KIWI, font_name=FONT_BODY)
    _add_text(slide, MARGIN, Inches(1.95), Inches(7), Inches(1.6),
              "Vi glæder os til at høre fra jer.",
              font_size=44, bold=True, color=WHITE, font_name=FONT_DISPLAY)

    # 2 kontakt-kort
    kam = ctx["team"]["kam"]
    rm = ctx["team"]["rm"]
    contacts = [
        ("DIN KEY ACCOUNT MANAGER", kam, KIWI),
        ("DIN RESOURCE MANAGER", rm, RASPBERRY),
    ]
    card_w = Inches(5.5)
    card_x = SLIDE_W - card_w - MARGIN
    card_y = Inches(2.0)
    card_h = Inches(2.3)
    for i, (role_label, person, accent) in enumerate(contacts):
        y = card_y + i * (card_h + Inches(0.1))
        _add_rect(slide, card_x, y, Emu(40000), card_h, accent)
        _add_rect(slide, card_x + Emu(40000), y,
                  card_w - Emu(40000), card_h,
                  RGBColor(0x25, 0x25, 0x70))
        _add_text(slide, card_x + Inches(0.3), y + Inches(0.25),
                  card_w - Inches(0.5), Inches(0.3),
                  role_label, font_size=10, bold=True, color=accent,
                  font_name=FONT_BODY)
        _add_text(slide, card_x + Inches(0.3), y + Inches(0.6),
                  card_w - Inches(0.5), Inches(0.5),
                  person.get("name", "[Navn]"), font_size=22, bold=True,
                  color=WHITE, font_name=FONT_DISPLAY)
        _add_text(slide, card_x + Inches(0.3), y + Inches(1.1),
                  card_w - Inches(0.5), Inches(0.3),
                  f"{person.get('title', '')} · Epico DK",
                  font_size=11,
                  color=RGBColor(0xAA, 0xAA, 0xBB),
                  font_name=FONT_BODY)
        _add_text(slide, card_x + Inches(0.3), y + Inches(1.5),
                  card_w - Inches(0.5), Inches(0.7),
                  f"T: {person.get('phone', '')}\nM: {person.get('email', '')}",
                  font_size=11, color=WHITE, font_name=FONT_BODY)


# ============================================================
# MAIN RENDERER
# ============================================================

def _strip_md(text: str) -> str:
    """Fjern **fed** markup — PPTX bruger ren tekst."""
    return (text or "").replace("**", "")


def _slide_from_library(prs, s: Dict[str, Any]) -> None:
    """
    Render én bibliotek-slide til PPTX ud fra dens layout.
    Spejler library_slide.html.j2 så HTML og PPTX matcher.
    """
    dark = s.get("variant") == "dark"
    red = s.get("variant") == "red"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BLACK_CURRANT if dark else (RED if red else RAW_SILK))
    _add_slide_header(slide, s.get("section_tag") or s.get("title", ""), dark=dark or red)

    heading_color = WHITE if (dark or red) else BLACK_CURRANT
    body_color = RGBColor(0xC8, 0xC8, 0xD8) if (dark or red) else GREY
    label_color = KIWI if (dark or red) else RASPBERRY

    y = Inches(1.15)

    if s.get("eyebrow"):
        _add_text(slide, MARGIN, y, Inches(9), Inches(0.3),
                  s["eyebrow"].upper(), font_size=10, bold=True,
                  color=label_color, font_name=FONT_BODY)
        y += Inches(0.42)

    if s.get("heading"):
        heading = _strip_md(s["heading"])
        size = 44 if len(heading) < 60 else 34
        _add_text(slide, MARGIN, y, SLIDE_W - (MARGIN * 2), Inches(1.0),
                  heading, font_size=size, bold=True,
                  color=heading_color, font_name=FONT_DISPLAY)
        y += Inches(0.95 if size == 44 else 1.15)

    if s.get("subheading"):
        _add_text(slide, MARGIN, y, Inches(10.5), Inches(0.6),
                  s["subheading"], font_size=15,
                  color=LIGHT_GREY if not (dark or red) else RGBColor(0xAA, 0xAA, 0xBB),
                  font_name=FONT_BODY)
        y += Inches(0.7)

    layout = s.get("layout", "bullets")
    stats = s.get("stats", [])
    bullets = s.get("bullets", [])
    cards = s.get("cards", [])
    content_w = SLIDE_W - (MARGIN * 2)

    # ---------- STATS HERO ----------
    if layout == "stats-hero" and stats:
        n = max(1, len(stats))
        col_w = int((content_w - Emu(20000) * (n - 1)) / n)
        for i, stat in enumerate(stats):
            x = MARGIN + i * (col_w + Emu(20000))
            accent = [KIWI, RASPBERRY, BLUEBERRY][i % 3]
            _add_rect(slide, x, y, Emu(25000), Inches(2.2), accent)
            _add_text(slide, x + Inches(0.22), y + Inches(0.25),
                      col_w - Inches(0.3), Inches(0.8),
                      stat.get("value", ""), font_size=34, bold=True,
                      color=WHITE if dark else BLACK_CURRANT, font_name=FONT_DISPLAY)
            _add_text(slide, x + Inches(0.22), y + Inches(1.15),
                      col_w - Inches(0.3), Inches(0.9),
                      stat.get("label", ""), font_size=10,
                      color=body_color, font_name=FONT_BODY)

    # ---------- STATS + BULLETS ----------
    elif layout == "stats-plus-bullets":
        if stats:
            n = max(1, len(stats))
            col_w = int((content_w - Emu(20000) * (n - 1)) / n)
            for i, stat in enumerate(stats):
                x = MARGIN + i * (col_w + Emu(20000))
                _add_rect(slide, x, y, col_w, Inches(1.5), WHITE if not dark else RGBColor(0x25, 0x25, 0x70))
                _add_rect(slide, x, y, Emu(25000), Inches(1.5), RED if not dark else KIWI)
                _add_text(slide, x + Inches(0.3), y + Inches(0.2),
                          col_w - Inches(0.5), Inches(0.6),
                          stat.get("value", ""), font_size=32, bold=True,
                          color=RASPBERRY if not dark else KIWI, font_name=FONT_DISPLAY)
                _add_text(slide, x + Inches(0.3), y + Inches(0.85),
                          col_w - Inches(0.5), Inches(0.5),
                          stat.get("label", ""), font_size=10,
                          color=body_color, font_name=FONT_BODY)
            y += Inches(1.75)
        _render_bullet_columns(slide, bullets, y, content_w, body_color, label_color)

    # ---------- CARDS ----------
    elif layout in ("cards-3", "cards-4", "cards-6", "competence-grid"):
        cols = {"cards-3": 3, "cards-4": 4, "cards-6": 3, "competence-grid": 4}[layout]
        rows = -(-len(cards) // cols) if cards else 1
        col_w = int((content_w - Emu(20000) * (cols - 1)) / cols)
        avail_h = SLIDE_H - y - Inches(0.7)
        card_h = int((avail_h - Emu(20000) * (rows - 1)) / max(1, rows))

        for i, card in enumerate(cards):
            cx = MARGIN + (i % cols) * (col_w + Emu(20000))
            cy = y + (i // cols) * (card_h + Emu(20000))
            _add_rect(slide, cx, cy, col_w, card_h, WHITE)
            _add_rect(slide, cx, cy, col_w, Emu(30000), RED)

            _add_text(slide, cx + Inches(0.25), cy + Inches(0.22),
                      col_w - Inches(0.45), Inches(0.5),
                      card.get("title", ""), font_size=15, bold=True,
                      color=BLACK_CURRANT, font_name=FONT_DISPLAY)

            inner_y = cy + Inches(0.75)
            if card.get("body"):
                _add_text(slide, cx + Inches(0.25), inner_y,
                          col_w - Inches(0.45), Inches(0.8),
                          card["body"], font_size=10, color=GREY, font_name=FONT_BODY)
                inner_y += Inches(0.85)

            for b in card.get("bullets", [])[:6]:
                _add_text(slide, cx + Inches(0.25), inner_y,
                          col_w - Inches(0.45), Inches(0.35),
                          f"· {b}", font_size=9, color=GREY, font_name=FONT_BODY)
                inner_y += Inches(0.32)

    # ---------- SERVICE DETAIL ----------
    elif layout == "service-detail":
        col_w = int((content_w - Emu(40000)) / 3)
        col_h = SLIDE_H - y - Inches(0.7)

        # Kolonne 1 — første kort
        if cards:
            _add_rect(slide, MARGIN, y, col_w, col_h, WHITE)
            _add_text(slide, MARGIN + Inches(0.3), y + Inches(0.25),
                      col_w - Inches(0.5), Inches(0.3),
                      cards[0].get("title", "").upper(), font_size=10, bold=True,
                      color=RASPBERRY, font_name=FONT_BODY)
            by = y + Inches(0.72)
            for b in cards[0].get("bullets", [])[:6]:
                _add_text(slide, MARGIN + Inches(0.3), by,
                          col_w - Inches(0.5), Inches(0.55),
                          f"· {b}", font_size=10, color=GREY, font_name=FONT_BODY)
                by += Inches(0.52)

        # Kolonne 2 — nøgletal (mørk)
        c2x = MARGIN + col_w + Emu(20000)
        if stats:
            _add_rect(slide, c2x, y, col_w, col_h, BLACK_CURRANT)
            _add_text(slide, c2x + Inches(0.3), y + Inches(0.25),
                      col_w - Inches(0.5), Inches(0.3),
                      "NØGLETAL", font_size=10, bold=True, color=KIWI, font_name=FONT_BODY)
            sy = y + Inches(0.8)
            for stat in stats[:4]:
                _add_text(slide, c2x + Inches(0.3), sy,
                          col_w - Inches(0.5), Inches(0.5),
                          stat.get("value", ""), font_size=26, bold=True,
                          color=KIWI, font_name=FONT_DISPLAY)
                _add_text(slide, c2x + Inches(0.3), sy + Inches(0.42),
                          col_w - Inches(0.5), Inches(0.35),
                          stat.get("label", ""), font_size=9,
                          color=RGBColor(0xAA, 0xAA, 0xBB), font_name=FONT_BODY)
                sy += Inches(0.85)

        # Kolonne 3 — andet kort
        c3x = MARGIN + (col_w * 2) + Emu(40000)
        if len(cards) > 1:
            _add_rect(slide, c3x, y, col_w, col_h, WHITE)
            _add_text(slide, c3x + Inches(0.3), y + Inches(0.25),
                      col_w - Inches(0.5), Inches(0.3),
                      cards[1].get("title", "").upper(), font_size=10, bold=True,
                      color=RASPBERRY, font_name=FONT_BODY)
            by = y + Inches(0.72)
            for b in cards[1].get("bullets", [])[:6]:
                _add_text(slide, c3x + Inches(0.3), by,
                          col_w - Inches(0.5), Inches(0.55),
                          f"· {b}", font_size=10, color=GREY, font_name=FONT_BODY)
                by += Inches(0.52)

    # ---------- TWO COL ----------
    elif layout == "two-col":
        half = int((content_w - Inches(0.8)) / 2)
        if s.get("body"):
            _add_text(slide, MARGIN, y, half, Inches(3.0),
                      s["body"], font_size=15, color=body_color, font_name=FONT_BODY)
        bx = MARGIN + half + Inches(0.8)
        by = y
        for b in bullets[:8]:
            _add_rect(slide, bx, by + Inches(0.1), Inches(0.18), Emu(28000),
                      RASPBERRY if not dark else KIWI)
            _add_text(slide, bx + Inches(0.35), by,
                      half - Inches(0.35), Inches(0.5),
                      b, font_size=13, color=body_color, font_name=FONT_BODY)
            by += Inches(0.48)

    # ---------- TEXT HERO ----------
    elif layout == "text-hero":
        _add_text(slide, MARGIN, y + Inches(0.5), content_w, Inches(2.5),
                  s.get("body", ""), font_size=32,
                  color=heading_color, font_name=FONT_DISPLAY)

    # ---------- BULLETS (default) ----------
    else:
        _render_bullet_columns(slide, bullets, y, content_w, body_color, label_color)

    if s.get("footnote"):
        _add_text(slide, MARGIN, SLIDE_H - Inches(0.95),
                  content_w, Inches(0.3),
                  s["footnote"], font_size=9, color=LIGHT_GREY, font_name=FONT_BODY)

    _add_footer(slide, s.get("title", ""), dark=dark or red)


def _render_bullet_columns(slide, bullets, y, content_w, body_color, accent_color):
    """Render bullets — én kolonne hvis få, to hvis mange."""
    if not bullets:
        return
    two_col = len(bullets) > 5
    col_w = int((content_w - Inches(0.8)) / 2) if two_col else content_w
    per_col = -(-len(bullets) // 2) if two_col else len(bullets)

    for i, b in enumerate(bullets):
        col = i // per_col if two_col else 0
        row = i % per_col if two_col else i
        bx = MARGIN + col * (col_w + Inches(0.8))
        by = y + row * Inches(0.52)
        _add_rect(slide, bx, by + Inches(0.11), Inches(0.2), Emu(28000), accent_color)
        _add_text(slide, bx + Inches(0.38), by,
                  col_w - Inches(0.4), Inches(0.5),
                  b, font_size=13, color=body_color, font_name=FONT_BODY)


def render_pptx(
    client_name: str,
    analysis: Dict[str, Any],
    meeting: Optional[Dict[str, str]] = None,
    team: Optional[Dict[str, Dict[str, str]]] = None,
    pitch_length: str = "medium",
    services: Optional[List[str]] = None,
    stakeholder: Optional[str] = None,
    excluded_slide_ids: Optional[List[str]] = None,
) -> bytes:
    """
    Generér et komplet pitch deck som .pptx.

    Samme to-delte struktur som HTML-versionen:
      DEL 1: Kunde-slides fra AI-analysen
      DEL 2: Epico-slides fra slide_library/, filtreret på længde + services
    """
    meeting = meeting or {}
    team = team or {}

    ctx = {
        "client": {"name": client_name},
        "meeting": {
            "date": meeting.get("date") or "—",
            "city": meeting.get("city") or "—",
            "contact_person": meeting.get("contact_person") or "—",
        },
        "team": {
            "kam": {
                "name": (team.get("kam") or {}).get("name") or "[Navn]",
                "title": (team.get("kam") or {}).get("title") or "Senior Key Account Manager",
                "phone": (team.get("kam") or {}).get("phone") or "+45 00 00 00 00",
                "email": (team.get("kam") or {}).get("email") or "[navn]@epico.dk",
            },
            "rm": {
                "name": (team.get("rm") or {}).get("name") or "[Navn]",
                "title": (team.get("rm") or {}).get("title") or "Resource Manager",
                "phone": (team.get("rm") or {}).get("phone") or "+45 00 00 00 00",
                "email": (team.get("rm") or {}).get("email") or "[navn]@epico.dk",
            },
        },
        "research_facts": analysis.get("research_facts", []),
        "value_mappings": analysis.get("value_mappings", []),
        "next_steps": analysis.get("next_steps", []),
        "case": analysis.get("case_recommendation", {}),
        "industry_tag": analysis.get("industry_tag", "branchen"),
        "headlines": analysis.get("slide_headlines") or {},
    }

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # DEL 1 — kunde-slides
    _slide_cover(prs, ctx)
    _slide_agenda(prs, ctx)
    _slide_research(prs, ctx)
    _slide_mapping(prs, ctx)

    # DEL 2 — Epico-slides fra biblioteket
    library = select_slides(
        pitch_length=pitch_length,
        services=services,
        stakeholder=stakeholder,
        excluded_ids=excluded_slide_ids,
    )
    if library:
        if pitch_length != "short":
            _slide_divider(prs, "02", "Dette er Epico.")
        for lib_slide in library:
            _slide_from_library(prs, lib_slide.to_dict())

    # DEL 3 — afslutning
    if ctx["case"] and ctx["case"].get("headline"):
        _slide_case(prs, ctx)
    _slide_next_steps(prs, ctx)
    _slide_contact(prs, ctx)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
